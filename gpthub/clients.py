from __future__ import annotations

from asyncio import to_thread
from io import BytesIO

from aioboto3 import Session
from aiobotocore.client import AioBaseClient
from config_fastapi import Config
from duckduckgo_search import DDGS
from pptx import Presentation
from webparser.http import HttpParser

from .constants import CACHE_TTL_WEB_PAGE, CACHE_TTL_WEB_SEARCH, WEB_SEARCH_MAX_RESULTS
from .structures import File, WebPage
from .utils import acached


class WebParserClient:

    def __init__(self):
        self.session = HttpParser()

    def __repr__(self) -> str:
        return type(self).__name__

    @property
    def client(self) -> HttpParser:
        return self.session

    @acached(ttl=CACHE_TTL_WEB_PAGE)
    async def get(self, url: str) -> WebPage:
        content = await self.client.get(url)
        return WebPage.create({
            "url": url,
            "text": content.text[:3000] if content.text else "",
        })

    async def fetch_bytes(self, url: str) -> WebPage:
        content = await self.client.fetch_bytes(url)
        return WebPage.create({
            "url": url,
            "content": content.content,
            "content_type": content.content_type or "application/octet-stream",
        })


class WebSearchClient:

    def __repr__(self) -> str:
        return type(self).__name__

    @property
    def client(self) -> DDGS:
        return DDGS()

    @acached(ttl=CACHE_TTL_WEB_SEARCH)
    async def search(self, query: str, max_results: int = WEB_SEARCH_MAX_RESULTS) -> list[WebPage]:
        results = await to_thread(self.client.text, query, max_results=max_results)
        return [
            WebPage.create({
                "url": r.get("href", ""),
                "title": r.get("title", ""),
                "text": r.get("body", "")[:500],
            })
            for r in results
        ]


class PPTXClient:

    async def build(self, title: str, slides: list) -> bytes:
        def render() -> bytes:
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = title
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data.title
                body = slide.shapes.placeholders[1].text_frame
                body.text = slide_data.bullets[0] if slide_data.bullets else ""
                for bullet in slide_data.bullets[1:]:
                    body.add_paragraph().text = bullet
            buffer = BytesIO()
            prs.save(buffer)
            return buffer.getvalue()
        return await to_thread(render)


class S3Client:
    config = Config(section="api.s3")

    def __init__(self) -> None:
        self.session = Session(
            aws_access_key_id=self.config.user,
            aws_secret_access_key=self.config.password,
            region_name=self.config.get("region"),
        )

    @property
    def client(self) -> AioBaseClient:
        return self.session.client("s3", endpoint_url=self.config.get("url"), region_name=self.config.get("region"))

    async def upload(self, file: File) -> File:
        async with self.client as s3cli:
            await s3cli.put_object(
                Bucket=self.config.bucket,
                Key=file.with_extension,
                Body=file.content,
                ContentType=file.extension.mime,
            )
        return file.set_url(f"{self.config.get('public_url')}/{self.config.bucket}/{file.with_extension}")
